"""
PowerPoint Report Generator — converts a GetHeard research report dict into a
branded .pptx slide deck ready for boardroom presentation.

Slides produced:
  1.  Title
  2.  Executive Summary
  3.  At-a-Glance (key stat + sentiment bar)
  4.  Respondent Personas (one card per persona, up to 4)
  5.  Key Themes (visual frequency bars)
  6.  Emotional Journey (stage list)
  7.  Pain Points (severity-coded)
  8.  Positive Highlights
  9.  Recommendations (by priority)
  10. Notable Quotes
  11. Research Gaps & Next Steps
  12. Thank You / Contact
"""

import io
import logging
from datetime import datetime
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

logger = logging.getLogger(__name__)

# ── Brand colours ──────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1e, 0x3c, 0x72)
NAVY_D = RGBColor(0x0f, 0x22, 0x44)


def _parse_hex(hex_str: str) -> RGBColor:
    """Parse '#rrggbb' → RGBColor. Falls back to NAVY on any error."""
    try:
        h = hex_str.lstrip("#")
        if len(h) == 6:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        pass
    return NAVY


def _accent(branding: Optional[dict]) -> RGBColor:
    """Return the primary brand accent colour (or default navy)."""
    if branding and branding.get("brand_color"):
        return _parse_hex(branding["brand_color"])
    return NAVY
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x6b, 0x72, 0x80)
LGRAY  = RGBColor(0xf0, 0xf4, 0xf8)
GREEN  = RGBColor(0x22, 0xc5, 0x5e)
AMBER  = RGBColor(0xf5, 0x9e, 0x0b)
RED    = RGBColor(0xef, 0x44, 0x44)
BLUE   = RGBColor(0x3b, 0x82, 0xf6)

SENTIMENT_COLORS = {"positive": GREEN, "neutral": AMBER, "negative": RED, "mixed": AMBER}
PRIORITY_COLORS  = {"high": RED, "medium": AMBER, "low": GREEN}
SEV_COLORS       = {"high": RED, "medium": AMBER, "low": GREEN}

SLIDE_W = Cm(33.87)   # 16:9 widescreen
SLIDE_H = Cm(19.05)

LANG_NAMES = {
    "en": "English", "hi": "Hindi", "id": "Indonesian", "fil": "Filipino",
    "th": "Thai", "vi": "Vietnamese", "ko": "Korean", "ja": "Japanese", "zh": "Mandarin",
}


# ── Helpers ────────────────────────────────────────────────────────────────────

def _blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def _rect(slide, l, t, w, h, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def _txt(slide, text: str, l, t, w, h, size: int, bold=False, color=WHITE,
         align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def _header_band(slide, title: str, subtitle: str = "", accent: RGBColor = None):
    """Dark band at top of slide with title. Uses accent colour."""
    band_color = accent or NAVY_D
    _rect(slide, 0, 0, SLIDE_W, Cm(3.6), band_color)
    _txt(slide, title, Cm(1.2), Cm(0.5), Cm(28), Cm(1.8),
         size=22, bold=True, color=WHITE)
    if subtitle:
        _txt(slide, subtitle, Cm(1.2), Cm(2.0), Cm(28), Cm(1.2),
             size=11, color=RGBColor(0xb0, 0xc4, 0xde))


def _bullet_box(slide, items: list, l, t, w, h, size=11, color=RGBColor(0x1f, 0x29, 0x37)):
    """Multi-line bullet text box."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(4)


def _pill(slide, text: str, l, t, color: RGBColor, text_color=WHITE, size=9):
    w = max(Cm(2.0), Cm(len(text) * 0.22 + 0.6))
    _rect(slide, l, t, w, Cm(0.55), color)
    _txt(slide, text.upper(), l + Cm(0.15), t + Cm(0.04), w, Cm(0.5),
         size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    return w


def _freq_bar(slide, label: str, pct: float, row: int, color: RGBColor):
    """Horizontal frequency bar row."""
    y = Cm(4.0) + row * Cm(1.3)
    max_bar_w = Cm(18)
    bar_w = max(Cm(0.3), max_bar_w * (pct / 100))
    _txt(slide, label, Cm(1.2), y, Cm(8), Cm(1.0), size=11, color=RGBColor(0x1f, 0x29, 0x37))
    _rect(slide, Cm(9.5), y + Cm(0.15), bar_w, Cm(0.6), color)
    _txt(slide, f"{pct:.0f}%", Cm(9.5) + bar_w + Cm(0.2), y, Cm(2), Cm(1.0),
         size=10, color=GRAY, bold=True)


# ── Slides ──────────────────────────────────────────────────────────────────────

def _slide_title(prs, report: dict, branding: Optional[dict] = None):
    slide = _blank_slide(prs)
    brand_color = _accent(branding)
    # Full-bleed background using brand colour (darkened via NAVY_D fallback)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY_D)
    # Accent stripe using brand colour
    _rect(slide, 0, Cm(12), SLIDE_W, Cm(0.25), brand_color)
    # Brand/client name if provided
    brand_name = branding.get("brand_name", "") if branding else ""
    brand_label = f"🎙 GetHeard" + (f"  ·  {brand_name}" if brand_name else "")
    _txt(slide, brand_label, Cm(1.5), Cm(1.2), Cm(20), Cm(1.5),
         size=16, bold=True, color=RGBColor(0x93, 0xc5, 0xfd))
    # Project name
    _txt(slide, report.get("project_name", "Research Report"),
         Cm(1.5), Cm(4.5), Cm(27), Cm(4),
         size=36, bold=True, color=WHITE, wrap=True)
    # Objective
    obj = report.get("objective", "")
    if obj:
        _txt(slide, obj[:140], Cm(1.5), Cm(9.0), Cm(26), Cm(2.5),
             size=14, color=RGBColor(0xb0, 0xc4, 0xde), wrap=True)
    # Meta strip
    meta = [
        f"👥 {report.get('total_transcripts', 0)} Respondents",
        f"🔬 {(report.get('research_type', 'CX')).upper()}",
        f"🌐 {', '.join(LANG_NAMES.get(l, l) for l in report.get('languages', []))}",
        f"📅 {datetime.fromisoformat(report.get('generated_at', datetime.now().isoformat())).strftime('%d %b %Y')}",
    ]
    _txt(slide, "  |  ".join(meta), Cm(1.5), Cm(12.5), Cm(28), Cm(1.2),
         size=11, color=RGBColor(0x7d, 0xa0, 0xd0))
    # Confidential footer
    conf_text = "Confidential — Prepared by GetHeard"
    if brand_name:
        conf_text = f"Confidential — Prepared for {brand_name} by GetHeard"
    _txt(slide, conf_text, Cm(1.5), Cm(17.5), Cm(28), Cm(1),
         size=9, color=RGBColor(0x4a, 0x6a, 0x9a), italic=True)


def _slide_exec_summary(prs, report: dict, branding: Optional[dict] = None):
    slide = _blank_slide(prs)
    _header_band(slide, "Executive Summary", accent=_accent(branding))
    _rect(slide, 0, Cm(3.6), SLIDE_W, SLIDE_H - Cm(3.6), LGRAY)
    summary = report.get("executive_summary", "No summary available.")
    # Split into paragraphs — show first ~600 chars across 2 columns
    paras = [p.strip() for p in summary.split("\n") if p.strip()]
    col1 = "\n\n".join(paras[:max(1, len(paras)//2)])
    col2 = "\n\n".join(paras[max(1, len(paras)//2):]) if len(paras) > 1 else ""

    _txt(slide, col1, Cm(1.2), Cm(4.0), Cm(14.5), Cm(14),
         size=11, color=RGBColor(0x1f, 0x29, 0x37), wrap=True)
    if col2:
        _txt(slide, col2, Cm(17), Cm(4.0), Cm(14.5), Cm(14),
             size=11, color=RGBColor(0x1f, 0x29, 0x37), wrap=True)

    # Confidence note
    conf = report.get("confidence_notes", "")
    if conf:
        _txt(slide, f"ℹ {conf[:200]}", Cm(1.2), Cm(17.0), Cm(30), Cm(1.5),
             size=9, color=GRAY, italic=True, wrap=True)


def _slide_at_a_glance(prs, report: dict, branding: Optional[dict] = None):
    slide = _blank_slide(prs)
    _header_band(slide, "At a Glance", "Key metrics from this study", accent=_accent(branding))
    _rect(slide, 0, Cm(3.6), SLIDE_W, SLIDE_H - Cm(3.6), LGRAY)

    # Big key stat
    key_stat = report.get("key_stat", "")
    if key_stat:
        _rect(slide, Cm(1.2), Cm(4.0), Cm(31), Cm(2.6), WHITE)
        _txt(slide, f"💡 {key_stat}", Cm(1.5), Cm(4.2), Cm(30), Cm(2.2),
             size=14, bold=True, color=NAVY, wrap=True)

    # 4 stat boxes
    sent = report.get("sentiment_overview", {})
    themes = report.get("key_themes", [])
    recs   = report.get("recommendations", [])
    stats = [
        (str(report.get("total_transcripts", 0)), "Respondents"),
        (str(len(themes)),                         "Key Themes"),
        (str(len([r for r in recs if r.get("priority") == "high"])), "High-Priority Actions"),
        ((sent.get("overall") or "—").capitalize(), "Overall Sentiment"),
    ]
    for i, (val, lbl) in enumerate(stats):
        x = Cm(1.2) + i * Cm(8)
        _rect(slide, x, Cm(7.2), Cm(7.4), Cm(3.2), WHITE)
        _txt(slide, val, x + Cm(0.3), Cm(7.5), Cm(6.8), Cm(1.8),
             size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _txt(slide, lbl, x + Cm(0.3), Cm(9.1), Cm(6.8), Cm(1),
             size=10, color=GRAY, align=PP_ALIGN.CENTER)

    # Sentiment bar
    _txt(slide, "Sentiment Breakdown", Cm(1.2), Cm(11.0), Cm(20), Cm(1),
         size=12, bold=True, color=NAVY)
    pos = sent.get("positive_pct", 0)
    neu = sent.get("neutral_pct", 0)
    neg = sent.get("negative_pct", 0)
    total_w = Cm(30)
    bar_h = Cm(0.8)
    bar_y = Cm(12.0)
    _rect(slide, Cm(1.2), bar_y, total_w * pos / 100 if pos else 0, bar_h, GREEN)
    _rect(slide, Cm(1.2) + total_w * pos / 100, bar_y, total_w * neu / 100 if neu else 0, bar_h, AMBER)
    _rect(slide, Cm(1.2) + total_w * (pos + neu) / 100, bar_y, total_w * neg / 100 if neg else 0, bar_h, RED)
    legend = f"  ■ Positive {pos}%    ■ Neutral {neu}%    ■ Negative {neg}%"
    _txt(slide, legend, Cm(1.2), Cm(13.0), Cm(28), Cm(1), size=10, color=GRAY)

    narrative = sent.get("sentiment_narrative", "")
    if narrative:
        _txt(slide, narrative[:200], Cm(1.2), Cm(14.2), Cm(30), Cm(3),
             size=11, color=RGBColor(0x1f, 0x29, 0x37), wrap=True, italic=True)


def _slide_personas(prs, report: dict, branding: Optional[dict] = None):
    personas = report.get("personas", [])
    if not personas:
        return
    slide = _blank_slide(prs)
    _header_band(slide, "Respondent Personas", "Behavioral archetypes derived from the interview data", accent=_accent(branding))
    _rect(slide, 0, Cm(3.6), SLIDE_W, SLIDE_H - Cm(3.6), LGRAY)

    avatars = ["👤", "👥", "🧑‍💼", "👩‍💻"]
    cols = min(len(personas), 4)
    card_w = Cm(30.5 / cols)

    for i, p in enumerate(personas[:4]):
        x = Cm(1.2) + i * (card_w + Cm(0.3))
        _rect(slide, x, Cm(4.0), card_w, Cm(14.2), WHITE)
        # Avatar + name
        _txt(slide, avatars[i % 4], x + Cm(0.3), Cm(4.2), Cm(2), Cm(1.2), size=24)
        _txt(slide, f"{p.get('percentage', '?')}%", x + card_w - Cm(1.8), Cm(4.3),
             Cm(1.5), Cm(0.8), size=10, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
        _txt(slide, p.get("name", "Persona"), x + Cm(0.3), Cm(5.5),
             card_w - Cm(0.6), Cm(1.4), size=12, bold=True, color=NAVY, wrap=True)
        _txt(slide, p.get("description", "")[:180], x + Cm(0.3), Cm(7.0),
             card_w - Cm(0.6), Cm(3.0), size=9, color=RGBColor(0x37, 0x41, 0x51), wrap=True)
        # Traits
        traits = ", ".join(p.get("characteristics", [])[:3])
        _txt(slide, f"Traits: {traits}", x + Cm(0.3), Cm(10.2),
             card_w - Cm(0.6), Cm(1.0), size=9, color=NAVY, bold=True)
        # What they need
        need = p.get("what_they_need", "")[:120]
        _txt(slide, f"Needs: {need}", x + Cm(0.3), Cm(11.2),
             card_w - Cm(0.6), Cm(2.0), size=9, color=RGBColor(0x37, 0x41, 0x51), wrap=True)
        # Quote
        quote = p.get("key_quote", "")
        if quote:
            _rect(slide, x + Cm(0.2), Cm(13.4), card_w - Cm(0.4), Cm(4.5),
                  RGBColor(0xe8, 0xee, 0xf7))
            _txt(slide, f'"{quote[:160]}"', x + Cm(0.4), Cm(13.6),
                 card_w - Cm(0.8), Cm(4.1), size=9, color=NAVY, italic=True, wrap=True)


def _slide_themes(prs, report: dict, branding: Optional[dict] = None):
    themes = report.get("key_themes", [])
    if not themes:
        return
    slide = _blank_slide(prs)
    _header_band(slide, "Key Themes", f"Top {len(themes)} themes identified across interviews", accent=_accent(branding))
    _rect(slide, 0, Cm(3.6), SLIDE_W, SLIDE_H - Cm(3.6), LGRAY)

    max_freq = max((t.get("frequency", 1) for t in themes), default=1)
    for i, t in enumerate(themes[:10]):
        y = Cm(4.0) + i * Cm(1.45)
        freq = t.get("frequency", 0)
        pct  = t.get("frequency_pct", 0)
        sent = t.get("sentiment", "neutral")
        color = SENTIMENT_COLORS.get(sent, AMBER)

        # Theme name
        _txt(slide, t.get("theme", ""), Cm(1.2), y, Cm(10), Cm(1.1),
             size=11, color=RGBColor(0x1f, 0x29, 0x37), bold=True)
        # Bar
        bar_w = max(Cm(0.3), Cm(18) * freq / max_freq)
        _rect(slide, Cm(11.5), y + Cm(0.15), bar_w, Cm(0.65), color)
        _txt(slide, f"{freq} ({pct}%)", Cm(11.5) + bar_w + Cm(0.2), y,
             Cm(3), Cm(1.0), size=9, color=GRAY, bold=True)
        # Sentiment pill
        _pill(slide, sent, Cm(29), y + Cm(0.1), color, size=8)


def _slide_emotional_journey(prs, report: dict, branding: Optional[dict] = None):
    journey = report.get("emotional_journey", [])
    if not journey:
        return
    slide = _blank_slide(prs)
    _header_band(slide, "Emotional Journey", "How respondents felt across the interview stages", accent=_accent(branding))
    _rect(slide, 0, Cm(3.6), SLIDE_W, SLIDE_H - Cm(3.6), LGRAY)

    emotion_icons = {
        "curious": "🤔", "frustrated": "😤", "confused": "😕", "satisfied": "😊",
        "hopeful": "🌟", "anxious": "😰", "resigned": "😔", "happy": "😄",
        "neutral": "😐", "angry": "😠", "excited": "🤩", "relieved": "😌",
        "disappointed": "😞", "engaged": "👂", "open": "🙂",
    }

    for i, s in enumerate(journey[:7]):
        y = Cm(4.0) + i * Cm(2.1)
        valence = s.get("valence_score", 5)
        color = GREEN if valence >= 7 else (RED if valence <= 3 else AMBER)
        emotion = s.get("dominant_emotion", "neutral").lower()
        icon = emotion_icons.get(emotion, "💬")

        _rect(slide, Cm(1.2), y, Cm(30.5), Cm(1.8), WHITE)
        # Valence bar on left
        bar_h = Cm(1.8 * valence / 10)
        _rect(slide, Cm(1.2), y + Cm(1.8) - bar_h, Cm(0.4), bar_h, color)
        # Icon + stage name
        _txt(slide, icon, Cm(1.9), y + Cm(0.1), Cm(1.2), Cm(1), size=16)
        _txt(slide, s.get("stage", ""), Cm(3.2), y + Cm(0.1), Cm(10), Cm(0.85),
             size=11, bold=True, color=NAVY)
        _txt(slide, f"{emotion.capitalize()} · Valence {valence}/10",
             Cm(3.2), y + Cm(0.9), Cm(8), Cm(0.75), size=9, color=GRAY)
        desc = s.get("description", "")[:200]
        _txt(slide, desc, Cm(13.5), y + Cm(0.1), Cm(18), Cm(1.6),
             size=10, color=RGBColor(0x37, 0x41, 0x51), wrap=True)


def _slide_pain_points(prs, report: dict, branding: Optional[dict] = None):
    pains = report.get("pain_points", [])
    if not pains:
        return
    slide = _blank_slide(prs)
    _header_band(slide, "Pain Points", "Friction and frustrations identified across respondents", accent=_accent(branding))
    _rect(slide, 0, Cm(3.6), SLIDE_W, SLIDE_H - Cm(3.6), LGRAY)

    for i, p in enumerate(pains[:7]):
        y = Cm(4.0) + i * Cm(2.1)
        sev = p.get("severity", "medium")
        color = SEV_COLORS.get(sev, AMBER)
        _rect(slide, Cm(1.2), y, Cm(0.5), Cm(1.8), color)  # severity stripe
        _rect(slide, Cm(1.9), y, Cm(30), Cm(1.8), WHITE)
        _txt(slide, p.get("pain_point", ""), Cm(2.2), y + Cm(0.1),
             Cm(20), Cm(0.9), size=11, bold=True, color=RGBColor(0x1f, 0x29, 0x37))
        _pill(slide, sev, Cm(23), y + Cm(0.15), color, size=8)
        impact = p.get("business_impact", "")
        if impact:
            _txt(slide, f"💼 {impact[:120]}", Cm(2.2), y + Cm(0.85),
                 Cm(20), Cm(0.6), size=9, color=NAVY)
        example = p.get("example", "")
        if example:
            _txt(slide, f'"{example[:120]}"', Cm(2.2), y + Cm(1.2),
                 Cm(28), Cm(0.7), size=9, color=GRAY, italic=True)


def _slide_recommendations(prs, report: dict, branding: Optional[dict] = None):
    recs = report.get("recommendations", [])
    if not recs:
        return
    slide = _blank_slide(prs)
    _header_band(slide, "Recommendations", "Prioritised actions based on research findings", accent=_accent(branding))
    _rect(slide, 0, Cm(3.6), SLIDE_W, SLIDE_H - Cm(3.6), LGRAY)

    # Sort by priority
    order = {"high": 0, "medium": 1, "low": 2}
    sorted_recs = sorted(recs, key=lambda r: order.get(r.get("priority", "medium"), 1))

    for i, rec in enumerate(sorted_recs[:7]):
        y = Cm(4.0) + i * Cm(2.1)
        pri = rec.get("priority", "medium")
        color = PRIORITY_COLORS.get(pri, AMBER)
        _rect(slide, Cm(1.2), y, Cm(0.5), Cm(1.8), color)
        _rect(slide, Cm(1.9), y, Cm(30), Cm(1.8), WHITE)
        _txt(slide, rec.get("recommendation", ""), Cm(2.2), y + Cm(0.05),
             Cm(22), Cm(0.9), size=11, bold=True, color=RGBColor(0x1f, 0x29, 0x37), wrap=True)
        _pill(slide, pri, Cm(25), y + Cm(0.15), color, size=8)
        owner = rec.get("who_owns_it", "")
        if owner:
            _txt(slide, f"Owner: {owner}", Cm(27), y + Cm(0.15), Cm(4.5), Cm(0.7),
                 size=8, color=GRAY)
        impact = rec.get("expected_impact", "")[:120]
        if impact:
            _txt(slide, f"→ {impact}", Cm(2.2), y + Cm(0.95),
                 Cm(28), Cm(0.65), size=9, color=NAVY)


def _slide_opportunity_matrix(prs, report: dict, branding: Optional[dict] = None):
    opps = report.get("opportunity_matrix", [])
    if not opps:
        return
    slide = _blank_slide(prs)
    _header_band(slide, "Opportunity Matrix", "Impact vs effort — where to focus next", accent=_accent(branding))
    _rect(slide, 0, Cm(3.6), SLIDE_W, SLIDE_H - Cm(3.6), LGRAY)

    cat_colors = {"quick_win": GREEN, "strategic": BLUE, "fill_in": GRAY, "backburner": AMBER}
    cat_labels = {
        "quick_win": "QUICK WIN",
        "strategic": "STRATEGIC",
        "fill_in": "FILL-IN",
        "backburner": "BACKBURNER",
    }

    for i, o in enumerate(opps[:8]):
        y = Cm(4.0) + i * Cm(1.85)
        cat = o.get("category", "fill_in")
        color = cat_colors.get(cat, GRAY)
        _rect(slide, Cm(1.2), y, Cm(30.5), Cm(1.6), WHITE)
        _pill(slide, cat_labels.get(cat, cat), Cm(1.4), y + Cm(0.15), color, size=8)
        _txt(slide, o.get("recommendation", "")[:120], Cm(7.0), y + Cm(0.1),
             Cm(19), Cm(1.0), size=10, color=RGBColor(0x1f, 0x29, 0x37), bold=True, wrap=True)
        _txt(slide, f"Impact {o.get('impact_score',5)}/10 · Effort {o.get('effort_score',5)}/10",
             Cm(26.5), y + Cm(0.1), Cm(5), Cm(0.7), size=9, color=GRAY, align=PP_ALIGN.RIGHT)
        metric = o.get("business_metric", "")
        if metric:
            _txt(slide, f"→ {metric.replace('_',' ').title()}", Cm(7.0), y + Cm(1.0),
                 Cm(19), Cm(0.6), size=8, color=NAVY)


def _slide_quotes(prs, report: dict, branding: Optional[dict] = None):
    quotes = report.get("notable_quotes", [])
    if not quotes:
        return
    slide = _blank_slide(prs)
    _header_band(slide, "Notable Quotes", "Verbatim voices from respondents", accent=_accent(branding))
    _rect(slide, 0, Cm(3.6), SLIDE_W, SLIDE_H - Cm(3.6), LGRAY)

    sent_stripe = {"positive": GREEN, "negative": RED, "neutral": NAVY}
    cols, rows = 2, 3
    card_w = Cm(14.5)
    card_h = Cm(4.5)
    for i, q in enumerate(quotes[:cols * rows]):
        col = i % cols
        row = i // cols
        x = Cm(1.2) + col * Cm(16)
        y = Cm(4.0) + row * (card_h + Cm(0.3))
        sent = q.get("sentiment", "neutral")
        color = sent_stripe.get(sent, NAVY)
        _rect(slide, x, y, card_w, card_h, WHITE)
        _rect(slide, x, y, Cm(0.35), card_h, color)
        _txt(slide, f'"{q.get("quote","")[:200]}"', x + Cm(0.55), y + Cm(0.2),
             card_w - Cm(0.8), Cm(3.3), size=10, color=RGBColor(0x1f, 0x29, 0x37),
             italic=True, wrap=True)
        lang = LANG_NAMES.get(q.get("language", "en"), q.get("language", ""))
        context = q.get("context", "")[:60]
        _txt(slide, f"{lang} · {context}", x + Cm(0.55), y + card_h - Cm(0.8),
             card_w - Cm(0.8), Cm(0.6), size=8, color=GRAY)


def _slide_gaps(prs, report: dict, branding: Optional[dict] = None):
    gaps = report.get("research_gaps", [])
    if not gaps:
        return
    slide = _blank_slide(prs)
    _header_band(slide, "Research Gaps & Next Steps", "Areas for follow-up investigation", accent=_accent(branding))
    _rect(slide, 0, Cm(3.6), SLIDE_W, SLIDE_H - Cm(3.6), LGRAY)
    _bullet_box(slide, [f"{g}" for g in gaps[:8]],
                Cm(1.5), Cm(4.2), Cm(30), Cm(12),
                size=13, color=RGBColor(0x92, 0x40, 0x0e))


def _slide_thankyou(prs, report: dict, branding: Optional[dict] = None):
    slide = _blank_slide(prs)
    brand_color = _accent(branding)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY_D)
    _rect(slide, 0, Cm(8.5), SLIDE_W, Cm(0.2), brand_color)
    _txt(slide, "Thank You", Cm(1.5), Cm(5.0), Cm(28), Cm(4),
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    brand_name = branding.get("brand_name", "") if branding else ""
    prepared_line = f"Prepared for {brand_name} by GetHeard · getheard.space" if brand_name else "Prepared by GetHeard · getheard.space"
    _txt(slide, prepared_line,
         Cm(1.5), Cm(10.5), Cm(28), Cm(1.5),
         size=14, color=RGBColor(0x7d, 0xa0, 0xd0), align=PP_ALIGN.CENTER)
    _txt(slide, "Questions? Contact us to run follow-up research.",
         Cm(1.5), Cm(12.5), Cm(28), Cm(1.5),
         size=12, color=RGBColor(0xb0, 0xc4, 0xde), align=PP_ALIGN.CENTER, italic=True)


# ── Public API ─────────────────────────────────────────────────────────────────

def generate_pptx(report: dict, branding: Optional[dict] = None) -> bytes:
    """
    Convert a report dict into a branded .pptx and return as bytes.
    Pass `branding` dict with optional keys: brand_name, brand_color, logo_url.
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_title(prs, report, branding)
    _slide_exec_summary(prs, report, branding)
    _slide_at_a_glance(prs, report, branding)
    _slide_personas(prs, report, branding)
    _slide_themes(prs, report, branding)
    _slide_emotional_journey(prs, report, branding)
    _slide_pain_points(prs, report, branding)
    _slide_recommendations(prs, report, branding)
    _slide_opportunity_matrix(prs, report, branding)
    _slide_quotes(prs, report, branding)
    _slide_gaps(prs, report, branding)
    _slide_thankyou(prs, report, branding)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    logger.info(f"Generated PPTX for report {report.get('report_id')} — {prs.slides.__len__()} slides")
    return buf.read()
